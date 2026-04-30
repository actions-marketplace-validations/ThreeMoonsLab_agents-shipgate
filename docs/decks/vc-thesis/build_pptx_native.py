"""Native PPTX reconstruction — text editable directly in PowerPoint.

Each slide is built with native shapes and text boxes (editable, restyleable).
A handful of complex SVG diagrams (Slide 4 Gödel loop, Slide 5 FEP boundary)
and the syntax-highlighted editor on Slide 8 embed cropped image fragments
because rebuilding them as PowerPoint shapes would lose intent.

Output: build/deck-editable.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

DECK_DIR = Path(__file__).resolve().parent
BUILD = DECK_DIR / "build"
FRAGS = BUILD / "fragments"
LOGO = BUILD / "_logo-mark-light.png"
OUT = BUILD / "deck-editable.pptx"

# ---- Brand palette ---------------------------------------------------------
CREAM     = RGBColor(0xF5, 0xF0, 0xE5)
CREAM_2   = RGBColor(0xEC, 0xE5, 0xD5)
CREAM_3   = RGBColor(0xE0, 0xD7, 0xC0)
PAPER     = RGBColor(0xFB, 0xF7, 0xEE)
NAVY      = RGBColor(0x1A, 0x25, 0x30)
NAVY_2    = RGBColor(0x2A, 0x35, 0x40)
NAVY_DEEP = RGBColor(0x0E, 0x18, 0x20)
MUTED     = RGBColor(0x6B, 0x5F, 0x4D)
MUTED_2   = RGBColor(0x8B, 0x7E, 0x68)
MUTED_DARK= RGBColor(0xB5, 0xA9, 0x88)
RULE      = RGBColor(0xD4, 0xCC, 0xB8)
CRITICAL  = RGBColor(0xB8, 0x39, 0x2F)
CRITICAL_BG = RGBColor(0xF4, 0xD9, 0xD6)
HIGH      = RGBColor(0xC7, 0x6A, 0x2C)
HIGH_BG   = RGBColor(0xF4, 0xE2, 0xD0)
GOLD      = RGBColor(0xD4, 0xA8, 0x47)
ACCENT_GREEN = RGBColor(0xB5, 0xC9, 0x9B)

# Fonts — use Aptos with Calibri fallback (PowerPoint 365 default)
FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Consolas"

# ---- Slide geometry -------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_Y = Inches(0.55)


# ---- Low-level helpers -----------------------------------------------------

def _set_shape_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_shape_line(shape, color=None, width=None):
    line = shape.line
    if color is None:
        line.fill.background()
    else:
        line.color.rgb = color
        if width is not None:
            line.width = width


def set_background(slide, color):
    """Add a slide-sized rectangle as the back-most fill."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _set_shape_fill(rect, color)
    _set_shape_line(rect)
    # Send to back by inserting at index 0
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius is not None:
        # adjustment value 0..1 — keep small for subtle rounding
        rect.adjustments[0] = radius
    if fill is not None:
        _set_shape_fill(rect, fill)
    else:
        rect.fill.background()
    _set_shape_line(rect, line, line_w)
    return rect


def add_line(slide, x1, y1, x2, y2, color, width=Pt(1)):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    ln.line.color.rgb = color
    ln.line.width = width
    return ln


def add_text(slide, x, y, w, h, text,
             font=FONT_BODY, size=14, bold=False, italic=False,
             color=NAVY, align="left", anchor="top", letter_spacing=None,
             line_spacing=None):
    """Single-run text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    if anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}[align]
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs,
             align="left", anchor="top", line_spacing=None,
             default_font=FONT_BODY, default_size=14):
    """Multi-run text box. `runs` is a list of dicts:
       {text, size, bold, italic, color, font}.
       Use {"break": True} to start a new paragraph.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    if anchor == "middle": tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom": tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else: tf.vertical_anchor = MSO_ANCHOR.TOP

    paragraphs = [tf.paragraphs[0]]
    cur = paragraphs[0]
    cur.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}[align]
    if line_spacing: cur.line_spacing = line_spacing

    for r in runs:
        if r.get("break"):
            cur = tf.add_paragraph()
            cur.alignment = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}[align]
            if line_spacing: cur.line_spacing = line_spacing
            paragraphs.append(cur)
            continue
        run = cur.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", default_font)
        run.font.size = Pt(r.get("size", default_size))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", NAVY)
    return tb


def add_picture(slide, x, y, w, h, path):
    return slide.shapes.add_picture(str(path), x, y, w, h)


def add_footer(slide, slide_no, total=15, dark=False):
    fg = MUTED_DARK if dark else MUTED
    fg_strong = CREAM if dark else NAVY
    # Brand line
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  MARGIN_X, SLIDE_H - Inches(0.45),
                                  Inches(0.1), Inches(0.1))
    _set_shape_fill(dot, fg_strong)
    _set_shape_line(dot)
    add_text(slide, MARGIN_X + Inches(0.18), SLIDE_H - Inches(0.5),
             Inches(8), Inches(0.3),
             "Three Moons Lab · A working thesis · April 2026",
             font=FONT_BODY, size=10, bold=True, color=fg_strong)
    # Slide number
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2), SLIDE_H - Inches(0.5),
             Inches(2), Inches(0.3),
             f"{slide_no:02d} / {total:02d}",
             font=FONT_MONO, size=9, color=fg, align="right")


def add_kicker_headline(slide, kicker, headline_runs, dark=False,
                        headline_y=Inches(1.0), kicker_y=Inches(0.55),
                        headline_h=Inches(2.4)):
    kicker_color = MUTED_DARK if dark else MUTED_2
    add_text(slide, MARGIN_X, kicker_y, Inches(12), Inches(0.35),
             kicker, font=FONT_BODY, size=11, bold=True,
             color=kicker_color, letter_spacing=0.18)
    add_runs(slide, MARGIN_X, headline_y, Inches(12), headline_h,
             headline_runs, line_spacing=1.05,
             default_font=FONT_HEAD, default_size=44)


# ---- Slide builders --------------------------------------------------------

def build_slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    # Logo
    add_picture(s, MARGIN_X, Inches(1.1), Inches(1.0), Inches(1.0), LOGO)
    # Wordmark
    add_text(s, MARGIN_X, Inches(2.3), Inches(6), Inches(0.4),
             "T H R E E   M O O N S   L A B",
             font=FONT_BODY, size=11, bold=True, color=MUTED)
    # Headline
    add_runs(s, MARGIN_X, Inches(2.85), Inches(11), Inches(2.2),
             [{"text": "Release readiness", "size": 64, "bold": True, "color": NAVY},
              {"break": True},
              {"text": "for agentic systems.", "size": 64, "bold": True, "color": NAVY}],
             line_spacing=1.05, default_font=FONT_HEAD)
    # Subtitle italic
    add_text(s, MARGIN_X, Inches(5.05), Inches(10), Inches(0.5),
             "A working thesis — not a pitch.",
             font=FONT_BODY, size=20, italic=True, color=NAVY_2)
    # Author / date row
    add_runs(s, MARGIN_X, Inches(6.1), Inches(11), Inches(0.4),
             [{"text": "Wendy · pengfei@threemoonslab.com", "size": 10, "color": MUTED, "font": FONT_MONO},
              {"text": "      April 2026", "size": 10, "color": MUTED, "font": FONT_MONO},
              {"text": "      v0.1 — for discussion", "size": 10, "color": MUTED, "font": FONT_MONO}])
    add_footer(s, 1)


def build_slide_02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 1 · THE SHIFT",
                        [{"text": "Models that answer.", "size": 44, "bold": True, "color": NAVY},
                         {"break": True},
                         {"text": "Agents that ", "size": 44, "bold": True, "color": MUTED_2},
                         {"text": "act.", "size": 44, "bold": True, "color": NAVY}])

    # LEFT card — LLM call
    card_y = Inches(2.6)
    card_h = Inches(3.2)
    add_rect(s, MARGIN_X, card_y, Inches(6), card_h, fill=PAPER, line=RULE, line_w=Pt(0.75))
    add_text(s, MARGIN_X + Inches(0.3), card_y + Inches(0.2), Inches(5.4), Inches(0.3),
             "YESTERDAY — LLM CALL", font=FONT_BODY, size=10, bold=True, color=MUTED_2)
    add_text(s, MARGIN_X + Inches(0.3), card_y + Inches(0.5), Inches(5.4), Inches(0.5),
             "Input → Output", font=FONT_HEAD, size=22, bold=True, color=NAVY)
    # Mini diagram: prompt → model → text
    dy = card_y + Inches(1.1)
    box_w = Inches(1.3); box_h = Inches(0.5)
    add_rect(s, MARGIN_X + Inches(0.3), dy, box_w, box_h, fill=CREAM_2, line=RULE)
    add_text(s, MARGIN_X + Inches(0.3), dy, box_w, box_h, "prompt",
             font=FONT_MONO, size=11, color=NAVY, align="center", anchor="middle")
    add_line(s, MARGIN_X + Inches(1.65), dy + Inches(0.25),
             MARGIN_X + Inches(2.1), dy + Inches(0.25), MUTED, width=Pt(1.5))
    add_rect(s, MARGIN_X + Inches(2.1), dy, box_w, box_h, fill=CREAM_2, line=RULE)
    add_text(s, MARGIN_X + Inches(2.1), dy, box_w, box_h, "model",
             font=FONT_MONO, size=11, color=NAVY, align="center", anchor="middle")
    add_line(s, MARGIN_X + Inches(3.45), dy + Inches(0.25),
             MARGIN_X + Inches(3.9), dy + Inches(0.25), MUTED, width=Pt(1.5))
    add_text(s, MARGIN_X + Inches(3.9), dy, Inches(1.0), box_h, "text",
             font=FONT_MONO, size=11, color=NAVY, align="left", anchor="middle")

    add_runs(s, MARGIN_X + Inches(0.3), card_y + Inches(2.05), Inches(5.4), Inches(1.0),
             [{"text": "•  stateless", "size": 12, "color": NAVY_2},
              {"break": True},
              {"text": "•  no real-world side effects", "size": 12, "color": NAVY_2},
              {"break": True},
              {"text": "•  release risk: ", "size": 12, "color": NAVY_2},
              {"text": "\"is the answer wrong?\"", "size": 12, "italic": True, "color": MUTED}],
             line_spacing=1.6)

    # RIGHT card — Agent
    rx = MARGIN_X + Inches(6.4)
    add_rect(s, rx, card_y, Inches(6), card_h, fill=CRITICAL_BG, line=CRITICAL, line_w=Pt(0.75))
    add_text(s, rx + Inches(0.3), card_y + Inches(0.2), Inches(5.4), Inches(0.3),
             "TODAY — AGENT", font=FONT_BODY, size=10, bold=True, color=CRITICAL)
    add_text(s, rx + Inches(0.3), card_y + Inches(0.5), Inches(5.4), Inches(0.6),
             "Observe → Plan → Tool → Side effect → Memory",
             font=FONT_HEAD, size=18, bold=True, color=NAVY)

    # Mini diagram: agent loop ellipse → tool call → side effect → world
    dy = card_y + Inches(1.2)
    # Loop ellipse
    add_rect(s, rx + Inches(0.3), dy, Inches(1.5), Inches(0.7), fill=PAPER, line=NAVY)
    add_runs(s, rx + Inches(0.3), dy, Inches(1.5), Inches(0.7),
             [{"text": "observe", "size": 9, "color": NAVY, "font": FONT_MONO},
              {"break": True},
              {"text": "plan", "size": 9, "color": NAVY, "font": FONT_MONO},
              {"break": True},
              {"text": "memory", "size": 9, "color": NAVY, "font": FONT_MONO}],
             line_spacing=1.0, align="center", anchor="middle")
    add_line(s, rx + Inches(1.85), dy + Inches(0.35),
             rx + Inches(2.4), dy + Inches(0.35), CRITICAL, width=Pt(2))
    add_rect(s, rx + Inches(2.4), dy + Inches(0.1), Inches(1.5), Inches(0.5),
             fill=PAPER, line=CRITICAL)
    add_text(s, rx + Inches(2.4), dy + Inches(0.1), Inches(1.5), Inches(0.5),
             "side effect",
             font=FONT_MONO, size=10, bold=True, color=CRITICAL,
             align="center", anchor="middle")
    add_line(s, rx + Inches(3.95), dy + Inches(0.35),
             rx + Inches(4.55), dy + Inches(0.35), CRITICAL, width=Pt(2))
    add_text(s, rx + Inches(4.55), dy + Inches(0.1), Inches(0.9), Inches(0.5),
             "world", font=FONT_MONO, size=10, color=NAVY, align="left", anchor="middle")
    # tool call label above arrow
    add_text(s, rx + Inches(1.85), dy - Inches(0.05), Inches(0.55), Inches(0.3),
             "tool call", font=FONT_MONO, size=8, bold=True, color=CRITICAL, align="center")

    add_runs(s, rx + Inches(0.3), card_y + Inches(2.15), Inches(5.4), Inches(1.0),
             [{"text": "•  stateful, looping", "size": 12, "color": NAVY_2},
              {"break": True},
              {"text": "•  ", "size": 12, "color": NAVY_2},
              {"text": "real consequences", "size": 12, "bold": True, "color": NAVY},
              {"text": " — refunds, emails, PRs, deploys", "size": 12, "color": NAVY_2},
              {"break": True},
              {"text": "•  release risk: ", "size": 12, "color": NAVY_2},
              {"text": "\"did the agent do the wrong thing?\"", "size": 12, "italic": True, "color": CRITICAL}],
             line_spacing=1.6)

    # Bottom callout
    cy = Inches(6.0)
    add_rect(s, MARGIN_X, cy, Inches(12.1), Inches(0.85), fill=CREAM_2, line=NAVY, line_w=Pt(2))
    add_runs(s, MARGIN_X + Inches(0.3), cy + Inches(0.1), Inches(11.7), Inches(0.7),
             [{"text": "Once an agent can call tools, ", "size": 14, "color": NAVY},
              {"text": "every tool change becomes a release event.", "size": 14, "bold": True, "color": NAVY},
              {"text": "  The release process built for code does not map onto agents.", "size": 14, "color": NAVY_2}],
             anchor="middle", line_spacing=1.4)

    add_footer(s, 2)


def build_slide_03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 1 · THE SHIFT",
                        [{"text": "Agent Release Readiness", "size": 42, "bold": True, "color": NAVY},
                         {"break": True},
                         {"text": "is a new release decision.", "size": 42, "bold": True, "color": NAVY}])
    # Lede
    add_text(s, MARGIN_X, Inches(2.6), Inches(12), Inches(0.9),
             "Bounded assurance that a stochastic, open, tool-using system can enter a higher-permission "
             "environment — under a declared task scope, tool surface, permission boundary, and risk tier.",
             font=FONT_BODY, size=15, color=NAVY_2, line_spacing=1.4)

    # Three cards
    cards = [
        ("IT IS NOT — SOFTWARE TESTING",
         "Software testing assumes a deterministic code path. Agents make their action graph at runtime "
         "from goals, context, tools, and feedback."),
        ("IT IS NOT — LLM EVAL",
         "Eval scores measure input → output behavior on sampled tasks. They cannot answer whether this "
         "tool surface, in this environment, is safe to ship."),
        ("IT IS NOT — RUNTIME SRE",
         "SLOs, canaries, and observability fire during or after execution. Release readiness is the "
         "decision before we promote."),
    ]
    cw = Inches(3.95); ch = Inches(2.0); cy = Inches(3.85)
    for i, (label, body) in enumerate(cards):
        cx = MARGIN_X + i * (cw + Inches(0.1))
        add_rect(s, cx, cy, cw, ch, fill=PAPER, line=RULE)
        add_text(s, cx + Inches(0.25), cy + Inches(0.2), cw - Inches(0.5), Inches(0.3),
                 label, font=FONT_BODY, size=9, bold=True, color=MUTED_2)
        add_text(s, cx + Inches(0.25), cy + Inches(0.55), cw - Inches(0.5), ch - Inches(0.7),
                 body, font=FONT_BODY, size=12, color=NAVY_2, line_spacing=1.4)

    # Bottom dark callout
    cby = Inches(6.05)
    add_rect(s, MARGIN_X, cby, Inches(12.1), Inches(0.85), fill=NAVY)
    add_runs(s, MARGIN_X + Inches(0.3), cby + Inches(0.1), Inches(11.7), Inches(0.7),
             [{"text": "IT IS    ", "size": 11, "bold": True, "color": MUTED_DARK, "font": FONT_BODY},
              {"text": "An ", "size": 14, "color": CREAM},
              {"text": "evidence-based release decision", "size": 14, "bold": True, "color": GOLD},
              {"text": " over a stochastic, tool-using, state-mutating system — graded against a declared "
                       "operational envelope.", "size": 14, "color": CREAM}],
             anchor="middle", line_spacing=1.4)

    add_footer(s, 3)


def build_slide_04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, NAVY_DEEP)
    add_kicker_headline(s, "ACT 2 · FIRST PRINCIPLE №1",
                        [{"text": "A sufficiently capable agent", "size": 36, "bold": True, "color": CREAM},
                         {"break": True},
                         {"text": "cannot self-certify its own readiness.", "size": 36, "bold": True, "color": CREAM}],
                        dark=True, headline_h=Inches(2.0))

    # LEFT prose
    px = MARGIN_X; py = Inches(3.2); pw = Inches(6.5)
    add_runs(s, px, py, pw, Inches(3.5),
             [{"text": "Any system rich enough to express its own behavior contains statements about itself "
                       "it cannot prove from within.", "size": 15, "color": MUTED_DARK},
              {"break": True}, {"break": True},
              {"text": "For agents, those statements are about ", "size": 15, "color": MUTED_DARK},
              {"text": "side effects", "size": 15, "bold": True, "color": CREAM},
              {"text": ", ", "size": 15, "color": MUTED_DARK},
              {"text": "long-horizon consequence", "size": 15, "bold": True, "color": CREAM},
              {"text": ", and ", "size": 15, "color": MUTED_DARK},
              {"text": "prompt-injection susceptibility", "size": 15, "bold": True, "color": CREAM},
              {"text": ".", "size": 15, "color": MUTED_DARK},
              {"break": True}, {"break": True},
              {"text": "External assurance is not optional. It is structural.",
               "size": 15, "bold": True, "color": GOLD}],
             line_spacing=1.45)

    # RIGHT diagram (cropped image)
    add_picture(s, Inches(7.5), Inches(2.2), Inches(5.4), Inches(3.86),
                FRAGS / "godel-loop.png")

    add_footer(s, 4, dark=True)


def build_slide_05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, NAVY_DEEP)
    add_kicker_headline(s, "ACT 2 · FIRST PRINCIPLE №2",
                        [{"text": "Tool calls are where uncertainty", "size": 36, "bold": True, "color": CREAM},
                         {"break": True},
                         {"text": "escapes into the world.", "size": 36, "bold": True, "color": CREAM}],
                        dark=True, headline_h=Inches(2.0))

    # LEFT prose
    px = MARGIN_X; py = Inches(3.2); pw = Inches(6.5)
    add_runs(s, px, py, pw, Inches(3.5),
             [{"text": "Agents act to minimize prediction error under their generative model of reality — the ",
               "size": 15, "color": MUTED_DARK},
              {"text": "Free Energy Principle", "size": 15, "italic": True, "color": MUTED_DARK},
              {"text": " framing.", "size": 15, "color": MUTED_DARK},
              {"break": True}, {"break": True},
              {"text": "The ", "size": 15, "color": MUTED_DARK},
              {"text": "tool boundary", "size": 15, "bold": True, "color": CREAM},
              {"text": " is the only channel through which an agent's internal uncertainty becomes external side effect.",
               "size": 15, "color": MUTED_DARK},
              {"break": True}, {"break": True},
              {"text": "Release readiness = bounding free energy at the action boundary ",
               "size": 15, "bold": True, "color": GOLD},
              {"text": "before", "size": 15, "italic": True, "color": MUTED_DARK},
              {"text": " it propagates.", "size": 15, "bold": True, "color": GOLD}],
             line_spacing=1.45)

    # RIGHT diagram (cropped image)
    add_picture(s, Inches(7.3), Inches(2.5), Inches(5.5), Inches(3.3),
                FRAGS / "fep-boundary.png")

    add_footer(s, 5, dark=True)


def build_slide_06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 3 · OUR THESIS",
                        [{"text": "The evidence layer between", "size": 42, "bold": True, "color": NAVY},
                         {"break": True},
                         {"text": "agent dev and production action.", "size": 42, "bold": True, "color": NAVY}])

    # LEFT — stack
    sx = MARGIN_X; sy = Inches(2.8); sw = Inches(7.0)
    # ABOVE
    add_rect(s, sx, sy, sw, Inches(0.75), fill=PAPER, line=RULE)
    add_text(s, sx + Inches(0.25), sy + Inches(0.07), sw, Inches(0.25),
             "ABOVE", font=FONT_BODY, size=9, bold=True, color=MUTED_2)
    add_text(s, sx + Inches(0.25), sy + Inches(0.32), sw - Inches(0.5), Inches(0.4),
             "Agent frameworks · OpenAI Agents SDK · Anthropic · Google ADK · LangChain · CrewAI",
             font=FONT_BODY, size=12, color=NAVY_2)
    # MIDDLE — Three Moons
    sy2 = sy + Inches(0.85)
    add_rect(s, sx, sy2, sw, Inches(1.15), fill=NAVY)
    add_text(s, sx + Inches(0.25), sy2 + Inches(0.1), sw, Inches(0.25),
             "THREE MOONS LAB — WHAT'S MISSING",
             font=FONT_BODY, size=9, bold=True, color=GOLD)
    add_text(s, sx + Inches(0.25), sy2 + Inches(0.4), sw - Inches(0.5), Inches(0.4),
             "CI/CD + audit layer for agentic systems",
             font=FONT_HEAD, size=18, bold=True, color=CREAM)
    add_text(s, sx + Inches(0.25), sy2 + Inches(0.78), sw - Inches(0.5), Inches(0.3),
             "pre-release evidence · trace-based replay · runtime continuous readiness",
             font=FONT_BODY, size=11, color=MUTED_DARK)
    # BELOW
    sy3 = sy2 + Inches(1.25)
    add_rect(s, sx, sy3, sw, Inches(0.75), fill=PAPER, line=RULE)
    add_text(s, sx + Inches(0.25), sy3 + Inches(0.07), sw, Inches(0.25),
             "BELOW", font=FONT_BODY, size=9, bold=True, color=MUTED_2)
    add_text(s, sx + Inches(0.25), sy3 + Inches(0.32), sw - Inches(0.5), Inches(0.4),
             "Tool surfaces · MCP · OpenAPI · function tools · shell · computer use",
             font=FONT_BODY, size=12, color=NAVY_2)
    # Adjacency note
    add_text(s, sx, sy3 + Inches(0.95), sw, Inches(0.3),
             "Adjacent (not us): eval frameworks · runtime guardrails · LLM observability · MCP gateways.",
             font=FONT_BODY, size=10, italic=True, color=MUTED)

    # RIGHT — thesis card
    tx = sx + sw + Inches(0.4); ty = Inches(2.8); tw = Inches(4.7); th = Inches(3.6)
    add_rect(s, tx, ty, tw, th, fill=CREAM_2)
    # Left edge line
    add_rect(s, tx, ty, Inches(0.05), th, fill=NAVY)
    add_text(s, tx + Inches(0.3), ty + Inches(0.25), tw - Inches(0.5), Inches(0.3),
             "THESIS", font=FONT_BODY, size=10, bold=True, color=MUTED_2)
    add_runs(s, tx + Inches(0.3), ty + Inches(0.65), tw - Inches(0.5), Inches(2.0),
             [{"text": "Every production agent will need a ", "size": 16, "color": NAVY},
              {"text": "release-readiness record", "size": 16, "bold": True, "color": NAVY},
              {"text": " before it gets promoted — and a ", "size": 16, "color": NAVY},
              {"text": "trace-replayable evidence trail", "size": 16, "bold": True, "color": NAVY},
              {"text": " after.", "size": 16, "color": NAVY}],
             line_spacing=1.35)
    add_runs(s, tx + Inches(0.3), ty + Inches(2.4), tw - Inches(0.5), Inches(1.0),
             [{"text": "That record won't live inside the model. It won't live inside the framework. "
                       "It has to live in ", "size": 13, "color": NAVY_2},
              {"text": "independent infrastructure", "size": 13, "bold": True, "color": NAVY},
              {"text": ".", "size": 13, "color": NAVY_2}],
             line_spacing=1.4)

    add_footer(s, 6)


def build_slide_07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 3 · OUR WEDGE",
                        [{"text": "Tool-use is the right wedge.", "size": 44, "bold": True, "color": NAVY}])

    cards = [
        ("①", "Action boundary",
         "Tool call = the moment language becomes consequence. It's where every interesting risk "
         "crystallizes: side effect, scope, approval, idempotency, recoverability.",
         "The model becoming smarter doesn't change this boundary. It only makes it more active."),
        ("②", "Most formalizable",
         "Tool surfaces ship with structure: schemas, scopes, MCP annotations, OpenAPI specs, SDK "
         "function signatures. Static analysis bites — unlike \"is the agent's reasoning correct?\"",
         "Formalize what's crisp · annotate what's contextual · review what's ambiguous."),
        ("③", "Highest-leverage risk",
         "AppWorld, ToolEmu, AgentDojo, τ-bench, AgentHarm — the academic evidence converges: "
         "tool-use is where current agents fail, where attacks land, where damage compounds.",
         "High-stakes tools (refund, email, deploy, delete) need readiness, not a benchmark score."),
    ]
    cw = Inches(3.95); ch = Inches(3.4); cy = Inches(2.4)
    for i, (num, title, body, foot) in enumerate(cards):
        cx = MARGIN_X + i * (cw + Inches(0.1))
        add_rect(s, cx, cy, cw, ch, fill=PAPER, line=RULE)
        add_text(s, cx + Inches(0.3), cy + Inches(0.2), Inches(1), Inches(0.7),
                 num, font=FONT_HEAD, size=32, bold=True, color=CRITICAL)
        add_text(s, cx + Inches(0.3), cy + Inches(0.9), cw - Inches(0.6), Inches(0.5),
                 title, font=FONT_HEAD, size=18, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.3), cy + Inches(1.45), cw - Inches(0.6), Inches(1.4),
                 body, font=FONT_BODY, size=11, color=NAVY_2, line_spacing=1.4)
        add_text(s, cx + Inches(0.3), cy + Inches(2.85), cw - Inches(0.6), Inches(0.5),
                 foot, font=FONT_BODY, size=9, italic=True, color=MUTED, line_spacing=1.4)

    # Bottom dark callout
    cby = Inches(6.0)
    add_rect(s, MARGIN_X, cby, Inches(12.1), Inches(0.9), fill=NAVY)
    add_runs(s, MARGIN_X + Inches(0.3), cby + Inches(0.1), Inches(11.7), Inches(0.7),
             [{"text": "Wedge logic:  ", "size": 13, "bold": True, "color": GOLD},
              {"text": "The narrowest cut where the static check is meaningful, the risk is real, "
                       "the buyer is identifiable, and the evidence corpus compounds. ", "size": 13, "color": CREAM},
              {"text": "Tool-use clears all four.", "size": 13, "bold": True, "color": CREAM}],
             anchor="middle", line_spacing=1.4)

    add_footer(s, 7)


def build_slide_08(prs):
    """Slide 8 — DECLARED. Python source + YAML config side-by-side.
    Kept as full-bleed image because the syntax-highlighted editor view
    can't be cleanly reconstructed in PowerPoint shapes."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_picture(s, 0, 0, SLIDE_W, SLIDE_H, BUILD / "slide-08.png")


def build_slide_09(prs):
    """Slide 9 — DETECTED. The shipgate report on the declared tool surface.
    Kept as full-bleed image so the report card hierarchy stays intact."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_picture(s, 0, 0, SLIDE_W, SLIDE_H, BUILD / "slide-09.png")


def build_slide_10(prs):
    """Slide 10 — How the contract gets written (init → author → scan → iterate).
    Kept as full-bleed image because of the dense 4-card flow + 3-column footer."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_picture(s, 0, 0, SLIDE_W, SLIDE_H, BUILD / "slide-10.png")


def build_slide_11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 4 · THE PRODUCT PATH",
                        [{"text": "Beyond static — sandbox & trace.", "size": 42, "bold": True, "color": NAVY}])

    cards = [
        ("PHASE 2", "~6–12 months out", "Sandbox & simulation",
         "Turn the unknowns surfaced in Phase 1 into experimental evidence — without exposing production state.",
         ["Mocked tool execution & failure injection",
          "Prompt-injection harness on read-tools (web, email, docs)",
          "State-diff assertions for collateral damage",
          "Synthetic adversarial scenarios (ToolEmu / AppWorld lineage)"],
         "Output: pre-promotion stress test report. CI-attachable. Fails-loud on regressions."),
        ("PHASE 3", "~12–24 months out", "Trace, replay, runtime",
         "Turn one-time pre-release reports into a continuous readiness state — pulled from the agent's actual production behavior.",
         ["Trace ingestion: OpenAI Agents SDK, MCP events, custom hooks",
          "Replay bundles for incident forensics",
          "Regression detection across prompt / model / tool changes",
          "Runtime anomaly & blast-radius monitors"],
         "Output: living readiness state. Audit-grade. Connected to incident review."),
    ]
    cw = Inches(6.0); ch = Inches(3.95); cy = Inches(2.4)
    for i, (label, when, title, lede, bullets, foot) in enumerate(cards):
        cx = MARGIN_X + i * (cw + Inches(0.1))
        add_rect(s, cx, cy, cw, ch, fill=PAPER, line=RULE)
        add_text(s, cx + Inches(0.3), cy + Inches(0.18), Inches(2), Inches(0.3),
                 label, font=FONT_BODY, size=10, bold=True, color=MUTED_2)
        add_text(s, cx + cw - Inches(2.3), cy + Inches(0.18), Inches(2), Inches(0.3),
                 when, font=FONT_MONO, size=9, color=MUTED, align="right")
        add_text(s, cx + Inches(0.3), cy + Inches(0.5), cw - Inches(0.6), Inches(0.5),
                 title, font=FONT_HEAD, size=22, bold=True, color=NAVY)
        add_text(s, cx + Inches(0.3), cy + Inches(1.05), cw - Inches(0.6), Inches(0.8),
                 lede, font=FONT_BODY, size=11, color=NAVY_2, line_spacing=1.4)
        bullet_runs = []
        for j, b in enumerate(bullets):
            if j > 0: bullet_runs.append({"break": True})
            bullet_runs.append({"text": "•  " + b, "size": 11, "color": NAVY_2})
        add_runs(s, cx + Inches(0.3), cy + Inches(1.95), cw - Inches(0.6), Inches(1.5),
                 bullet_runs, line_spacing=1.5)
        add_line(s, cx + Inches(0.3), cy + ch - Inches(0.6),
                 cx + cw - Inches(0.3), cy + ch - Inches(0.6), RULE)
        add_text(s, cx + Inches(0.3), cy + ch - Inches(0.5), cw - Inches(0.6), Inches(0.4),
                 foot, font=FONT_BODY, size=10, italic=True, color=MUTED, line_spacing=1.4)

    add_text(s, MARGIN_X, Inches(6.55), Inches(12.1), Inches(0.4),
             "Phase 1 ships now. Phase 2 and 3 are deliberate land-and-expand, not a roadmap to be promised on Slide 9.",
             font=FONT_BODY, size=11, color=NAVY_2, align="center", italic=True)

    add_footer(s, 11)


def build_slide_12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 4 · WHY THIS IS A CATEGORY, NOT A FEATURE",
                        [{"text": "Three phases. One compounding", "size": 42, "bold": True, "color": NAVY},
                         {"break": True},
                         {"text": "evidence corpus.", "size": 42, "bold": True, "color": NAVY}])

    # LEFT — stacked layers
    sx = MARGIN_X; sw = Inches(7.6)
    # Phase 3 (top)
    p3y = Inches(2.6); p3h = Inches(1.05)
    add_rect(s, sx, p3y, sw, p3h, fill=CREAM_3, line=NAVY)
    add_text(s, sx + Inches(0.3), p3y + Inches(0.1), sw, Inches(0.25),
             "PHASE 3 · TRACE", font=FONT_BODY, size=10, bold=True, color=MUTED_2)
    add_text(s, sx + Inches(0.3), p3y + Inches(0.35), sw - Inches(0.5), Inches(0.3),
             "Production trace data", font=FONT_HEAD, size=16, bold=True, color=NAVY)
    add_text(s, sx + Inches(0.3), p3y + Inches(0.65), sw - Inches(0.5), Inches(0.4),
             "tool-call events · approval logs · replay bundles · regression deltas · incident forensics",
             font=FONT_MONO, size=10, color=NAVY_2)
    # Phase 2 (middle)
    p2y = p3y + p3h + Inches(0.05); p2h = Inches(0.95)
    add_rect(s, sx, p2y, sw, p2h, fill=CREAM_2, line=NAVY)
    add_text(s, sx + Inches(0.3), p2y + Inches(0.1), sw, Inches(0.25),
             "PHASE 2 · SANDBOX", font=FONT_BODY, size=10, bold=True, color=MUTED_2)
    add_text(s, sx + Inches(0.3), p2y + Inches(0.35), sw - Inches(0.5), Inches(0.3),
             "Failure-mode taxonomy", font=FONT_HEAD, size=16, bold=True, color=NAVY)
    add_text(s, sx + Inches(0.3), p2y + Inches(0.65), sw - Inches(0.5), Inches(0.3),
             "attack patterns · injection results · state-diff baselines · scenario library",
             font=FONT_MONO, size=10, color=NAVY_2)
    # Phase 1 (bottom)
    p1y = p2y + p2h + Inches(0.05); p1h = Inches(0.95)
    add_rect(s, sx, p1y, sw, p1h, fill=PAPER, line=NAVY)
    add_text(s, sx + Inches(0.3), p1y + Inches(0.1), sw, Inches(0.25),
             "PHASE 1 · STATIC", font=FONT_BODY, size=10, bold=True, color=MUTED_2)
    add_text(s, sx + Inches(0.3), p1y + Inches(0.35), sw - Inches(0.5), Inches(0.3),
             "Tool surface metadata", font=FONT_HEAD, size=16, bold=True, color=NAVY)
    add_text(s, sx + Inches(0.3), p1y + Inches(0.65), sw - Inches(0.5), Inches(0.3),
             "manifests · schemas · scopes · effect classes · approval flags",
             font=FONT_MONO, size=10, color=NAVY_2)
    # Compound arrow on right of the stack
    arrow_x = sx + sw + Inches(0.15)
    arrow_top_y = p3y + Inches(0.05)
    arrow_bottom_y = p1y + p1h - Inches(0.05)
    add_line(s, arrow_x, arrow_bottom_y, arrow_x, arrow_top_y + Inches(0.18),
             CRITICAL, width=Pt(2.5))
    # Arrow head triangle pointing up
    head = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                               arrow_x - Inches(0.12), arrow_top_y,
                               Inches(0.24), Inches(0.22))
    _set_shape_fill(head, CRITICAL)
    _set_shape_line(head)
    # "corpus compounds" caption below the stack — keeps it out of the right column
    add_text(s, sx, p1y + p1h + Inches(0.15), sw, Inches(0.3),
             "corpus compounds across phases →",
             font=FONT_BODY, size=11, bold=True, italic=True,
             color=CRITICAL, align="right")

    # RIGHT — explanation
    rx = sx + sw + Inches(1.2); rw = Inches(3.9)
    add_text(s, rx, Inches(2.6), rw, Inches(1.2),
             "Three phases are not three products. They are the same evidence corpus unfolding "
             "across three timescales.",
             font=FONT_BODY, size=14, color=NAVY_2, line_spacing=1.45)
    add_runs(s, rx, Inches(3.85), rw, Inches(1.0),
             [{"text": "Each user adds metadata, failure cases, and traces. The ", "size": 12, "color": NAVY_2},
              {"text": "failure taxonomy", "size": 12, "bold": True, "color": NAVY},
              {"text": ", ", "size": 12, "color": NAVY_2},
              {"text": "policy library", "size": 12, "bold": True, "color": NAVY},
              {"text": ", and ", "size": 12, "color": NAVY_2},
              {"text": "trace schema", "size": 12, "bold": True, "color": NAVY},
              {"text": " compound.", "size": 12, "color": NAVY_2}],
             line_spacing=1.4)
    # Anti-feature box
    afy = Inches(4.85)
    add_rect(s, rx, afy, rw, Inches(1.7), fill=NAVY)
    add_runs(s, rx + Inches(0.2), afy + Inches(0.15), rw - Inches(0.4), Inches(1.4),
             [{"text": "This is the anti-feature defense.", "size": 11, "bold": True, "color": GOLD},
              {"text": " A single GitHub Action lint cannot compound. A scanner backed by a growing "
                       "cross-organizational evidence corpus can.", "size": 11, "color": CREAM}],
             line_spacing=1.4)

    add_footer(s, 12)


def build_slide_13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 5 · WHERE I AM",
                        [{"text": "What I'm validating now.", "size": 44, "bold": True, "color": NAVY}])
    add_text(s, MARGIN_X, Inches(2.05), Inches(12), Inches(0.5),
             "Three open hypotheses. Each one converts (or kills) the company. The next 6–8 weeks are "
             "a validation loop, not a product sprint.",
             font=FONT_BODY, size=14, color=NAVY_2, line_spacing=1.4)

    rows = [
        ("H1",
         "Production agents have a recurring pre-release readiness workflow today — even if no one has named it.",
         "3–5 design partners running shipgate in real CI · 10+ release-blocking findings on real tool surfaces · repeatable trigger event."),
        ("H2",
         "The first owner is platform / AI infra engineering, not security/GRC. Security buys later, after evidence accumulates.",
         "Design-partner data on who triggers / triages findings · which team owns the CI gate · whether security review piggybacks on shipgate output."),
        ("H3",
         "Static + manifest checks are sufficient through Risk Tier 3 (reversible internal write). Tier 4+ requires sandbox + trace.",
         "Real findings on real tool surfaces, post-fix · false-positive rate on static checks · which tiers actually demand simulation evidence."),
    ]
    rh = Inches(1.05); ry0 = Inches(2.85)
    for i, (label, hyp, proof) in enumerate(rows):
        ry = ry0 + i * (rh + Inches(0.15))
        add_rect(s, MARGIN_X, ry, Inches(12.1), rh, fill=PAPER, line=RULE)
        add_text(s, MARGIN_X + Inches(0.25), ry + Inches(0.2), Inches(0.8), Inches(0.7),
                 label, font=FONT_HEAD, size=22, bold=True, color=CRITICAL, anchor="middle")
        add_text(s, MARGIN_X + Inches(1.1), ry + Inches(0.13), Inches(2.5), Inches(0.25),
                 "HYPOTHESIS", font=FONT_BODY, size=8, bold=True, color=MUTED_2)
        add_text(s, MARGIN_X + Inches(1.1), ry + Inches(0.38), Inches(5.3), Inches(0.65),
                 hyp, font=FONT_BODY, size=11, color=NAVY, line_spacing=1.3)
        add_text(s, MARGIN_X + Inches(6.6), ry + Inches(0.13), Inches(3), Inches(0.25),
                 "PROOF I'M SEEKING", font=FONT_BODY, size=8, bold=True, color=MUTED_2)
        add_text(s, MARGIN_X + Inches(6.6), ry + Inches(0.38), Inches(5.4), Inches(0.65),
                 proof, font=FONT_BODY, size=11, color=NAVY_2, line_spacing=1.3)

    add_text(s, MARGIN_X, Inches(6.4), Inches(12.1), Inches(0.4),
             "What I'm not claiming: PMF, runtime safety certification, or that this is foundation-model-lab-proof. "
             "Those are unknowns to be earned.",
             font=FONT_BODY, size=11, italic=True, color=MUTED, align="center")

    add_footer(s, 13)


def build_slide_14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "ACT 5 · 10-YEAR NORTH STAR",
                        [{"text": "Today: pre-release exam.", "size": 36, "bold": True, "color": NAVY},
                         {"break": True},
                         {"text": "Long-term: a record for every agent.", "size": 36, "bold": True, "color": NAVY}],
                        headline_h=Inches(2.0))

    # Timeline axis
    ty = Inches(4.3)
    add_line(s, MARGIN_X + Inches(0.5), ty, SLIDE_W - MARGIN_X - Inches(0.5), ty, NAVY, width=Pt(1.5))

    # Stages: (x, label, title, sub1, sub2, color, marker_text)
    stages = [
        (Inches(1.4),  "TODAY",            "Pre-release exam",  "static scanner ·",  "CI gate · SARIF report", NAVY,    "⊕"),
        (Inches(4.5),  "YEAR 2",           "Stress tests",      "sandbox · failure injection ·", "prompt-injection harness", NAVY, "⊗"),
        (Inches(7.6),  "YEAR 3–4",         "Vital signs",       "trace ingestion · replay ·", "runtime anomaly monitors", NAVY, "♥"),
        (Inches(10.7), "10-YEAR NORTH STAR", "Medical record", "across the lifetime of every agent —", "development, incidents, retirement", CRITICAL, "★"),
    ]
    for cx, label, title, sub1, sub2, color, marker in stages:
        # Marker symbol above
        add_text(s, cx - Inches(0.4), ty - Inches(1.0), Inches(0.8), Inches(0.5),
                 marker, font=FONT_HEAD, size=24, color=color, align="center")
        # Dot on axis
        dot = slide_oval(s, cx - Inches(0.07), ty - Inches(0.07), Inches(0.14), Inches(0.14))
        _set_shape_fill(dot, color)
        _set_shape_line(dot)
        # Label below
        add_text(s, cx - Inches(1.5), ty + Inches(0.2), Inches(3.0), Inches(0.3),
                 label, font=FONT_BODY, size=10, bold=True, color=color if color == CRITICAL else MUTED_2,
                 align="center")
        add_text(s, cx - Inches(1.5), ty + Inches(0.5), Inches(3.0), Inches(0.4),
                 title, font=FONT_HEAD, size=16, bold=True, color=NAVY, align="center")
        add_text(s, cx - Inches(1.5), ty + Inches(0.95), Inches(3.0), Inches(0.3),
                 sub1, font=FONT_BODY, size=10, color=NAVY_2, align="center")
        add_text(s, cx - Inches(1.5), ty + Inches(1.2), Inches(3.0), Inches(0.3),
                 sub2, font=FONT_BODY, size=10, color=NAVY_2, align="center")

    # Bottom callout — moved up to clear the footer
    cby = Inches(6.25)
    add_rect(s, MARGIN_X, cby, Inches(12.1), Inches(0.6), fill=CREAM_2)
    add_rect(s, MARGIN_X, cby, Inches(0.05), Inches(0.6), fill=NAVY)
    add_text(s, MARGIN_X + Inches(0.3), cby + Inches(0.04), Inches(11.7), Inches(0.5),
             "Today we build the first instrument. The compounding ambition is to make every production agent's "
             "release, incident, and behavior traceable and accountable across its life.",
             font=FONT_BODY, size=11, color=NAVY, anchor="middle", line_spacing=1.4)

    add_footer(s, 14)


def slide_oval(slide, x, y, w, h):
    return slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)


def build_slide_15(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(s, CREAM)
    add_kicker_headline(s, "CLOSING",
                        [{"text": "What I'm looking for.", "size": 44, "bold": True, "color": NAVY}])
    add_text(s, MARGIN_X, Inches(2.05), Inches(12), Inches(0.5),
             "This deck is not a fundraise. It's an invitation to think alongside us. "
             "Three concrete asks, in order of immediate value:",
             font=FONT_BODY, size=14, color=NAVY_2, line_spacing=1.4)

    cards = [
        ("ASK 01", "Sparring partners",
         "Founders, operators, researchers willing to push back on the thesis. Especially: people "
         "who think this is a feature, not a category. I want to be wrong fast.",
         "Best for: Prateek, AI-infra peers, security/GRC operators, MCP & framework authors.",
         "light"),
        ("ASK 02 — MOST VALUABLE NOW", "Design partners",
         "Teams shipping production agents with non-trivial tool surfaces — refunds, customer comms, "
         "code execution, internal data access. I want to scan, find real risk, watch what gets fixed, "
         "learn what their CI actually demands.",
         "Looking for: 3–5 partners over the next 6–8 weeks.",
         "dark"),
        ("ASK 03 — LATER", "Capital optionality",
         "Not raising today. When the design-partner loop converts the thesis to traction, I'd like "
         "the conversation to continue with people who already understood the worldview.",
         "Trigger: 3+ design partners using shipgate findings to gate releases.",
         "light"),
    ]
    cw = Inches(3.95); ch = Inches(3.4); cy = Inches(2.65)
    for i, (label, title, body, foot, theme) in enumerate(cards):
        cx = MARGIN_X + i * (cw + Inches(0.1))
        if theme == "dark":
            add_rect(s, cx, cy, cw, ch, fill=NAVY)
            label_color = GOLD
            title_color = CREAM
            body_color = MUTED_DARK
            foot_color = MUTED_DARK
        else:
            add_rect(s, cx, cy, cw, ch, fill=PAPER, line=RULE)
            label_color = MUTED_2
            title_color = NAVY
            body_color = NAVY_2
            foot_color = MUTED
        add_text(s, cx + Inches(0.3), cy + Inches(0.22), cw - Inches(0.6), Inches(0.3),
                 label, font=FONT_BODY, size=10, bold=True, color=label_color)
        add_text(s, cx + Inches(0.3), cy + Inches(0.6), cw - Inches(0.6), Inches(0.5),
                 title, font=FONT_HEAD, size=20, bold=True, color=title_color)
        add_text(s, cx + Inches(0.3), cy + Inches(1.2), cw - Inches(0.6), Inches(1.7),
                 body, font=FONT_BODY, size=11, color=body_color, line_spacing=1.4)
        add_text(s, cx + Inches(0.3), cy + ch - Inches(0.45), cw - Inches(0.6), Inches(0.35),
                 foot, font=FONT_BODY, size=9, italic=True, color=foot_color)

    # Bottom callout — moved up to clear the footer (text wraps to 2 lines)
    cby = Inches(6.2)
    cbh = Inches(0.65)
    add_rect(s, MARGIN_X, cby, Inches(12.1), cbh, fill=CREAM_2)
    add_rect(s, MARGIN_X, cby, Inches(0.05), cbh, fill=CRITICAL)
    add_runs(s, MARGIN_X + Inches(0.3), cby + Inches(0.08), Inches(11.7), cbh - Inches(0.16),
             [{"text": "Three Moons Lab is not building infrastructure to make agents smarter. "
                       "We're building infrastructure to make their entry into the world ",
               "size": 11, "italic": True, "color": NAVY},
              {"text": "accountable", "size": 11, "italic": True, "bold": True, "color": CRITICAL},
              {"text": ".", "size": 11, "italic": True, "color": NAVY}],
             anchor="middle", line_spacing=1.4)

    add_footer(s, 15)


# ---- Main -----------------------------------------------------------------

BUILDERS = [
    build_slide_01, build_slide_02, build_slide_03, build_slide_04, build_slide_05,
    build_slide_06, build_slide_07, build_slide_08, build_slide_09, build_slide_10,
    build_slide_11, build_slide_12, build_slide_13, build_slide_14, build_slide_15,
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for builder in BUILDERS:
        builder(prs)

    prs.save(OUT)
    size_kb = OUT.stat().st_size // 1024
    print(f"Wrote {OUT}  ({size_kb} KB)")


if __name__ == "__main__":
    main()
