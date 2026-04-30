"""Build deck.pptx — each slide is the rendered PNG full-bleed,
with a speaker note derived from the reading guide.

Visual fidelity: pixel-perfect (each slide is the rendered image).
Editability: structural (reorder, add/remove slides, edit speaker notes).
NOT editable: text on slides — to change copy, edit build_deck.py and re-run.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

BUILD_DIR = Path(__file__).resolve().parent / "build"
OUT = BUILD_DIR / "deck.pptx"

# Speaker notes per slide — derived from the reading guide so Wendy has prompts
# while presenting. Edit freely in PowerPoint without affecting visuals.
NOTES = {
    1:  "Subtitle: 'A working thesis — not a pitch.' "
        "Set this expectation up front so the room reads correctly: I'm not raising; "
        "I'm pressure-testing whether this can become a venture-scale infrastructure company.",
    2:  "Highest-leverage slide for category creation. "
        "The bottom callout — 'every tool change becomes a release event' — is the entire wedge in one line. "
        "If they read past this and don't update their model of agents, the rest of the deck won't land.",
    3:  "Three negative-space cards (it is not — software testing / LLM eval / runtime SRE) before the positive definition. "
        "Defending the category claim by what it isn't is more honest than a TAM slide.",
    4:  "Gödel — *why* this category must exist. External assurance is structural, not a 'nice to have'. "
        "This slide closes off two competitors before they're raised: foundation labs ('we'll add this') "
        "and agent self-certification ('the agent can check itself'). "
        "Skip this slide if the audience is allergic to philosophy.",
    5:  "Free Energy — *where* the wedge intervenes. "
        "Tool calls are the only channel through which internal uncertainty becomes external side effect. "
        "Pivots naturally into Slide 7's tool-use argument. "
        "Skip if you skipped Slide 4.",
    6:  "Stack diagram positions us between agent frameworks (above) and tool surfaces (below), "
        "distinct from adjacent infra (eval, observability, MCP gateways). "
        "Right side: thesis stated plainly. "
        "If they push on 'why not LangSmith / Patronus / Galileo', point at the adjacency footer.",
    7:  "Three reasons tool-use is the right cut. "
        "Academic backup (AppWorld, ToolEmu, AgentDojo, τ-bench) shows this isn't a hunch. "
        "Wedge logic in the dark callout: any miss and tool-use isn't the right cut.",
    8:  "DECLARED — what the team built. Python (refund_agent.py) on the left shows the actual "
        "OpenAI Agents SDK wiring. YAML (shipgate.yaml) on the right shows the release contract. "
        "Read prohibited_actions out loud — 'issue refund without approval'. Hold that thought.",
    9:  "DETECTED — what shipgate found. Same tool surface, scanned. BLOCKED. 2 critical, 14 high. "
        "First finding: stripe.create_refund has no approval policy declared. "
        "That's the contradiction with the previous slide's prohibited_actions list. "
        "Let the audience connect Slide 8 ↔ Slide 9 themselves.",
    10: "How the contract gets written. Anticipates the question 'wait, where does shipgate.yaml come from?' "
        "Four steps: scaffold (auto) → author (human) → scan (auto) → iterate (human + CI). "
        "Bottom row makes the auto-vs-human split explicit — scanner can detect tools, only humans can write intent. "
        "The Terraform analogy works well here if they push: prose policy → policy as code.",
    11: "Sandbox turns Phase 1 unknowns into experimental evidence. "
        "Trace turns one-time reports into a continuous readiness state. "
        "Footer line is deliberate humility: 'Phase 1 ships now. Phase 2 and 3 are land-and-expand.'",
    12: "Most important slide for category vs feature framing. "
        "Three phases = same evidence corpus unfolding across timescales. "
        "Failure taxonomy / policy library / trace schema all compound with every user. "
        "If they say 'this could be a feature in [LangSmith / GitHub / Snyk]', this is the slide that pushes back.",
    13: "Three open hypotheses, each paired with proof being sought. "
        "Closes with what I'm NOT claiming — PMF, runtime safety certification, foundation-model-lab proof. "
        "The 'what I'm not claiming' line is the strongest credibility signal in the deck.",
    14: "Healthcare metaphor placed deliberately late so it doesn't dominate framing. "
        "Today: first instrument. Long-term: medical record across every agent's life. "
        "Skip this slide if the audience is too transactional for vision.",
    15: "Three asks in order of immediate value: sparring partners → design partners (most valuable now) → capital optionality. "
        "Closing italic line is the only place the deck declares philosophy directly: "
        "'infrastructure to make their entry into the world accountable.' "
        "End on this beat.",
}


def main():
    prs = Presentation()
    # 16:9 widescreen, exactly 1920x1080 in EMU at 144 DPI
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for n in range(1, 16):
        slide = prs.slides.add_slide(blank_layout)
        png = BUILD_DIR / f"slide-{n:02d}.png"
        if not png.exists():
            raise FileNotFoundError(png)

        # Full-bleed image
        slide.shapes.add_picture(
            str(png),
            left=0, top=0,
            width=prs.slide_width, height=prs.slide_height,
        )

        # Speaker notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = NOTES.get(n, "")

    prs.save(OUT)
    size_kb = OUT.stat().st_size // 1024
    print(f"Wrote {OUT}  ({size_kb} KB, {size_kb / 1024:.1f} MB)")


if __name__ == "__main__":
    main()
